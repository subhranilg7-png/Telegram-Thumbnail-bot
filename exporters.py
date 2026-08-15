from pathlib import Path
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from renderer import W,H,fit_crop,render


def export_png(project,out):
    render(project,1).save(out,'PNG')


def export_pdf(project,out):
    img=render(project,2)
    c=canvas.Canvas(out,pagesize=(W,H))
    c.drawImage(ImageReader(img),0,0,width=W,height=H)
    c.save()


def _xy(v,total,slide_total): return Inches(v/total*slide_total)

def _text(slide,text,x,y,w,h,size,bold=True,align=PP_ALIGN.LEFT,color=(255,255,255)):
    s=slide.shapes.add_textbox(_xy(x,W,13.333333),_xy(y,H,7.5),_xy(w,W,13.333333),_xy(h,H,7.5))
    tf=s.text_frame; tf.clear(); tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; r.font.name='Poppins'; r.font.size=Pt(size*0.75); r.font.bold=bold; r.font.color.rgb=RGBColor(*color)
    return s

def _shape(slide,x,y,w,h,fill=(25,25,30),trans=0,line=None,radius=True):
    typ=MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s=slide.shapes.add_shape(typ,_xy(x,W,13.333333),_xy(y,H,7.5),_xy(w,W,13.333333),_xy(h,H,7.5))
    s.fill.solid(); s.fill.fore_color.rgb=RGBColor(*fill); s.fill.transparency=int(trans)
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=RGBColor(*line); s.line.width=Pt(1)
    return s

def _tmp(img,out,suffix):
    p=Path(out).with_name(Path(out).stem+suffix+'.png'); img.save(p,'PNG'); return str(p)


def export_pptx(project,out):
    prs=Presentation(); prs.slide_width=Inches(13.333333); prs.slide_height=Inches(7.5)
    slide=prs.slides.add_slide(prs.slide_layouts[6]); p=project['layout']; accent=tuple(project.get('accent_color') or (220,60,70))

    # Main background only: no baked-in title/synopsis/rating.
    bg=Image.open(project['files']['background']).convert('RGB')
    bg=fit_crop(bg,(W,H),tuple(p.get('background_focal',[.5,.5])),p.get('background_zoom',1.0))
    slide.shapes.add_picture(_tmp(bg,out,'_bg'),0,0,width=prs.slide_width,height=prs.slide_height)
    # dark readability overlay
    _shape(slide,0,0,W,H,(0,0,0),58,None,False)

    # Header
    _text(slide,project.get('watermark','@Ongoing_english_dub'),45,42,330,36,19,True,PP_ALIGN.LEFT,accent)
    _text(slide,'•  THE BEST POSSIBLE ONGOING ANIME EXPERIENCE',96,88,360,22,8,False,PP_ALIGN.LEFT,(245,245,245))

    # Decorative geometry
    for x,y,s in [(70,22,42),(128,12,52),(190,35,38),(250,18,44)]:
        sh=slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,_xy(x,W,13.333333),_xy(y,H,7.5),_xy(s,W,13.333333),_xy(s,H,7.5))
        sh.fill.background(); sh.line.color.rgb=RGBColor(*accent); sh.line.transparency=45; sh.line.width=Pt(1)

    # Top image
    t=p.get('top_image',{}); tx,ty,tw,th=[int(t.get(k,v)) for k,v in [('x',18),('y',142),('w',610),('h',300)]]
    top=Image.open(project['files']['top']).convert('RGB'); top=fit_crop(top,(tw,th),tuple(t.get('focal',[.5,.5])),t.get('zoom',1.0))
    slide.shapes.add_picture(_tmp(top,out,'_top'),_xy(tx,W,13.333333),_xy(ty,H,7.5),width=_xy(tw,W,13.333333),height=_xy(th,H,7.5))
    _shape(slide,tx,ty,tw,th,(0,0,0),100,accent,True)

    # Synopsis
    syn=p.get('synopsis',{}); sx,sy,sw=[int(syn.get(k,v)) for k,v in [('x',38),('y',474),('width',610)]]
    _text(slide,'✦  SYNOPSIS  ✦',sx,sy,sw,36,21,True,PP_ALIGN.CENTER,(255,255,255))
    _text(slide,project.get('synopsis',''),sx,sy+40,sw,145,int(syn.get('body_size',15)),False,PP_ALIGN.CENTER,(250,250,250))
    _text(slide,'RATING:',150,582,130,36,22,True,PP_ALIGN.LEFT,(255,255,255))
    rating=project.get('rating'); rt='N/A' if rating in (None,'') else f'{float(rating):.1f}/10'
    _text(slide,rt,270,582,150,36,22,True,PP_ALIGN.LEFT,accent)
    _text(slide,'The ratings are sourced from a trusted anime\ninformation website called AniList',150,625,330,42,9,False,PP_ALIGN.LEFT,(245,245,245))

    # Cards
    for key,default in [('card1',p.get('card1',{})),('card2',p.get('card2',{}))]:
        c=default; x,y,w,h=[int(c.get(k,v)) for k,v in [('x',690),('y',490),('w',128),('h',120)]]
        im=Image.open(project['files'][key]).convert('RGB'); im=fit_crop(im,(w,h),tuple(c.get('focal',[.5,.5])),c.get('zoom',1.0))
        slide.shapes.add_picture(_tmp(im,out,'_'+key),_xy(x,W,13.333333),_xy(y,H,7.5),width=_xy(w,W,13.333333),height=_xy(h,H,7.5))
        _shape(slide,x,y,w,h,(0,0,0),100,accent,True)

    # Season badge
    season=project.get('season_text')
    if season:
        _shape(slide,1000,390,185,38,(15,15,20),55,accent,True)
        _text(slide,season.upper(),1000,393,185,30,11,False,PP_ALIGN.CENTER,accent)

    # Glass info panel
    pan=p.get('info_panel',{}); x,y,w,h=[int(pan.get(k,v)) for k,v in [('x',825),('y',440),('w',430),('h',250)]]
    _shape(slide,x,y,w,h,(22,22,28),35,(120,120,125),True)
    _shape(slide,x+w-48,y+20,32,32,accent,0,None,False)
    _shape(slide,x+w-48,y+h-52,32,32,accent,0,None,False)
    title=project.get('title','Untitled'); subtitle=project.get('subtitle') or project.get('romaji_title') or ''
    _text(slide,title,x+25,y+24,w-50,58,27,True,PP_ALIGN.CENTER,(255,255,255))
    if subtitle and subtitle.lower()!=title.lower(): _text(slide,subtitle,x+30,y+84,w-60,28,12,True,PP_ALIGN.CENTER,accent)
    _text(slide,'SEASON '+(season.split(' ',1)[0] if season else 'N/A'),x+40,y+120,w-80,30,13,False,PP_ALIGN.CENTER,(255,255,255))
    genres=project.get('genres') or []
    yy=y+157
    for row in [genres[:3],genres[3:6]]:
        if not row: continue
        total=sum(max(70,int(len(g)*6.8+28)) for g in row)+10*(len(row)-1); xx=x+(w-total)/2
        for i,g in enumerate(row):
            pw=max(70,int(len(g)*6.8+28)); active=(yy==y+157 and i==0)
            _shape(slide,xx,yy,pw,30,accent if active else (25,25,30),0 if active else 55,accent if not active else None,True)
            _text(slide,g.upper(),xx,yy+2,pw,24,8,False,PP_ALIGN.CENTER,(255,255,255)); xx+=pw+10
        yy+=38
    _text(slide,project.get('watermark','@Ongoing_english_dub'),x+50,y+h-38,w-100,25,10,False,PP_ALIGN.CENTER,accent)

    prs.save(out)
    for tmp in Path(out).parent.glob(Path(out).stem+'_*.png'):
        try: tmp.unlink()
        except OSError: pass
